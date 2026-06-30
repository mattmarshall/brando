#!/usr/bin/env python3
"""brando office — a branded PowerPoint deck (also imports into Google Slides).

Generalized from the fastverk pptx tool: the slide construction is brand-neutral;
the brand supplies its palette, font, title/tagline, color swatches, and the
mark/wordmark images via a JSON config.

CLI: pptx_gen.py <config.json> <out.pptx> <wordmark_dark.png> <icon_dark.png> <icon_light.png>

config.json:
{
  "name": "fastverk",
  "tagline": "proven systems, built fast",
  "font": "Space Grotesk",
  "bg": "#15161A", "fg": "#ECE7DA", "accent": "#E0A33E", "tertiary": "#4A565A",
  "swatches": [["#15161A","15161A"], ["#ECE7DA","ECE7DA"],
               ["#E0A33E","E0A33E"], ["#4A565A","4A565A"]]
}
"""
import json
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SW, SH = Inches(13.333), Inches(7.5)


def _rgb(hexstr):
    h = hexstr.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _text(slide, text, left, top, width, height, size, color, font, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        r.font.name = font
    return tb


def _center_pic(slide, path, top, height):
    pic = slide.shapes.add_picture(path, 0, top, height=height)
    pic.left = int((SW - pic.width) / 2)
    return pic


def build(cfg, out, wordmark, icon_dark, icon_light):
    name = cfg["name"]
    tagline = cfg.get("tagline", "")
    font = cfg.get("font", "Helvetica")
    BG, FG = _rgb(cfg["bg"]), _rgb(cfg["fg"])
    ACCENT, SLATE = _rgb(cfg["accent"]), _rgb(cfg["tertiary"])
    swatches = cfg.get("swatches", [[cfg["bg"], cfg["bg"].lstrip("#")],
                                    [cfg["fg"], cfg["fg"].lstrip("#")],
                                    [cfg["accent"], cfg["accent"].lstrip("#")],
                                    [cfg["tertiary"], cfg["tertiary"].lstrip("#")]])

    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    blank = prs.slide_layouts[6]

    # title
    s = prs.slides.add_slide(blank); _bg(s, BG)
    _center_pic(s, wordmark, Inches(2.3), Inches(1.7))
    _text(s, tagline, 0, Inches(4.4), SW, Inches(0.6), 22, SLATE, font, align=PP_ALIGN.CENTER)
    _text(s, "Brand deck · generated from one parametric source", 0, Inches(5.05), SW,
          Inches(0.5), 14, SLATE, font, align=PP_ALIGN.CENTER)

    # the mark
    s = prs.slides.add_slide(blank); _bg(s, BG)
    _text(s, "The mark", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, font, bold=True)
    _text(s, "One parametric source\nGeometry by exact CSG\nOne source for SVG, icons, brandbook, decks",
          Inches(0.9), Inches(1.8), Inches(7), Inches(4), 20, FG, font)
    s.shapes.add_picture(icon_dark, Inches(9.2), Inches(2.0), height=Inches(3.2))

    # light & dark
    s = prs.slides.add_slide(blank); _bg(s, BG)
    _text(s, "Light & dark", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, font, bold=True)
    s.shapes.add_picture(icon_dark, Inches(3.0), Inches(2.2), height=Inches(3.0))
    s.shapes.add_picture(icon_light, Inches(7.8), Inches(2.2), height=Inches(3.0))
    _text(s, "dark", Inches(3.0), Inches(5.4), Inches(2.6), Inches(0.5), 16, SLATE, font, align=PP_ALIGN.CENTER)
    _text(s, "light", Inches(7.8), Inches(5.4), Inches(2.6), Inches(0.5), 16, SLATE, font, align=PP_ALIGN.CENTER)

    # color
    s = prs.slides.add_slide(blank); _bg(s, BG)
    _text(s, "Color", Inches(0.8), Inches(0.5), Inches(8), Inches(0.8), 30, FG, font, bold=True)
    for i, (col, label) in enumerate(swatches):
        x = Inches(0.9 + i * 3.0)
        sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.2), Inches(2.4), Inches(1.7))
        sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(col); sh.line.color.rgb = SLATE
        _text(s, label, x, Inches(4.0), Inches(2.4), Inches(0.5), 14, FG, font)

    # closing
    s = prs.slides.add_slide(blank); _bg(s, BG)
    _center_pic(s, icon_dark, Inches(1.8), Inches(2.6))
    _text(s, name, 0, Inches(5.0), SW, Inches(0.8), 28, FG, font, bold=True, align=PP_ALIGN.CENTER)

    prs.save(out)
    print("wrote", out)


def main():
    cfg = json.load(open(sys.argv[1]))
    build(cfg, sys.argv[2], sys.argv[3], sys.argv[4], sys.argv[5])


if __name__ == "__main__":
    main()
